from pptx import Presentation
import sys

def inspect_presentation(filepath):
    prs = Presentation(filepath)
    print(f"Presentation has {len(prs.slides)} slides.")
    
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        for j, shape in enumerate(slide.shapes):
            shape_type = shape.shape_type if hasattr(shape, 'shape_type') else 'Unknown'
            print(f"  Shape {j}: Type={shape_type}, Name='{shape.name}'")
            if shape.has_text_frame:
                text = shape.text.replace("\n", "\\n").replace("\v", "\\v").strip()
                print(f"    Text: {text[:50]}..." if len(text) > 50 else f"    Text: {text}")

if __name__ == '__main__':
    if len(sys.argv) > 1:
        inspect_presentation(sys.argv[1])
    else:
        print("Usage: python inspect_ppt.py <path_to_pptx>")
