from pptx import Presentation
from pptx.util import Inches, Pt
import sys

def fill_presentation(template_path, output_path):
    prs = Presentation(template_path)
    
    # Data Mapping
    content_map = {
        0: { # Slide 1: Title
            8: "Intelligent Orbital Debris Removal Planner",
            7: "Track 4: Sustainable Space Systems",
            5: "Team ta5abes"
        },
        1: { # Slide 2: Problem Statement
            "title": "ProbeLem Statement",
            "body": "• Over 36,500 objects >10cm threaten orbital sustainability.\n• 1 million+ fragments (>1cm) posing catastrophic collision risk.\n• Kessler Syndrome: A self-perpetuating collision cascade.\n• Fuel Bottleneck: Every m/s saved extends mission lifetime."
        },
        2: { # Slide 3: Proposed Solution
            6: "Autonomous Dual-Layer AI Architecture",
            "body": "• RL Core: PPO Agent optimizes multi-target collection sequencing.\n• RAG Advisory: Grounded decisions via NASA/ESA standard retrieval.\n• Mission Goal: Maximize debris cleared per gram of propellant."
        },
        3: { # Slide 4: System Architecture
            "body": "• Environment: 3D Keplerian simulator with real-world TLE data.\n• Strategy: Fuel-optimal sequencing via Reinforcement Learning.\n• Safety: Retrieval-Augmented Generation for operational compliance."
        },
        4: { # Slide 5: Technical Implementation
            "body": "• Simulation: Custom Gymnasium env using orbital mechanics (Hohmann).\n• RL: Stable-Baselines3 (PPO) for sampled trajectory optimization.\n• RAG: FAISS vector store + LangChain indexing orbital safety protocols."
        },
        5: { # Slide 6: Results
            "body": "• Delta-V Reduction: 86.3% lower propellant cost than random baseline.\n• Fuel Efficiency: 45% more efficient (m/s per target) via PPO sequencing.\n• Rigor: Validated over 1.7M simulation steps on high-precision 3D physics."
        },
        6: { # Slide 7: Impact
            "body": "• Ecological: Mitigates collision risk to maintain LEO usability.\n• Economic: Extends robotic removal spacecraft service life.\n• Social: Safeguards global satellite infrastructure (GPS, Weather)."
        },
        7: { # Slide 8: SDG Alignment
            "body": "• SDG 9: Innovation in orbital management technology.\n• SDG 12: Responsible orbital usage and waste reduction.\n• SDG 17: Alignment with international space safety standards (IADC)."
        },
        8: { # Slide 9: Scalability
            "body": "• Multi-Agent: Coordinating fleets of debris removal spacecraft.\n• Real-time: Live API integration with Space-Track and ESA catalogs.\n• Transferable: Adaptable to GEO and MEO orbital regimes."
        },
        9: { # Slide 10: Challenges
            "body": "• Computation: Scaling deep RL to 10,000+ uncatalogued fragments.\n• Sensor Noise: Modeling uncertainty in debris tracking data.\n• Dynamics: Incorporating orbital perturbations (Drag, J2 effect)."
        },
        10: { # Slide 11: Conclusion
            "body": "• Validated: Proof-of-concept for AI-driven orbital sustainability.\n• Efficient: Massive fuel savings (86%+) via deep autonomous planning.\n• Ready: Optimized model and code prepared for Phase 2 scaling."
        }
    }

    for i, slide in enumerate(prs.slides):
        if i not in content_map:
            continue
            
        data = content_map[i]
        
        # Fill specific shapes by ID if provided
        for shape_id, text in data.items():
            if isinstance(shape_id, int):
                try:
                    slide.shapes[shape_id].text = text
                except:
                    pass
        
        # Add body content if present
        if "body" in data:
            # Find an empty or placeholder shape
            # Based on inspection, Shapes 5, 6, 7 are often empty
            target_shape = None
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text.strip()
                    # If it's a content placeholder or empty
                    if t == "" or "Lorem" in t or "Core idea" in t:
                        target_shape = shape
                        break
            
            if target_shape:
                target_shape.text = data["body"]
                # Adjust font size if needed
                for p in target_shape.text_frame.paragraphs:
                    p.font.size = Pt(18)
            else:
                # Add new textbox if no placeholder found
                txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(4))
                tf = txBox.text_frame
                tf.text = data["body"]
                for p in tf.paragraphs:
                    p.font.size = Pt(18)

    prs.save(output_path)
    print(f"Saved to {output_path}")

if __name__ == '__main__':
    fill_presentation(sys.argv[1], sys.argv[2])
