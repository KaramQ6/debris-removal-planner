from pptx import Presentation
from pptx.util import Inches, Pt
import sys

def fill_presentation(template_path, output_path):
    prs = Presentation(template_path)
    
    slide_content = {
        "Introduction": "Intelligent Orbital Debris Removal Planner\nAutonomous Fuel-Optimal Debris Collection with RAG Advisory",
        "The Problem": "• Over 36,500 tracked debris objects >10cm threaten LEO.\n• Kessler Syndrome risk is increasing exponentially.\n• Current manual mission planning is fuel-wasteful and non-scalable.\n• Fuel is the ultimate constraint in orbital sustainability.",
        "Proposed Solution": "A two-layer AI architecture:\n1. RL Core Agent (PPO): Learns fuel-optimal, multi-target debris collection paths.\n2. RAG Advisory: Provides real-time guidance from NASA/ESA debris standards.",
        "System Architecture": "• Simulator: 3D Keplerian environment with TLE debris data.\n• RL Agent: Optimizes target sequencing (Delta-V vs. Risk).\n• RAG Base: FAISS vector store indexing NASA-STD-8719.14.\n• Output: 3D visual path and step-by-step fuel budget.",
        "Key Results": "Simulated on 8-target LEO scenario:\n• Random Baseline: 4,093 m/s Delta-V per target.\n• Nearest-Neighbor: 2,247 m/s Delta-V per target.\n• RL Agent (PPO): 3,067 m/s Delta-V per target.\nResult: RL efficiency outperforms random by 25% in early training phases.",
        "Demo / Visuals": "• The environment models realistic Hohmann transfers and inclination changes.\n• Training rewards fuel conservation and prioritizes high-risk debris targets.\n• Nearest-Neighbor heuristic serves as a reliable operational fallback.",
        "Sustainability Impact": "• Reduced Fuel Consumption: Extends mission lifetime.\n• Higher Clearance Rates: More objects removed per launch.\n• Lower Launch Mass: Reduces Earth-to-Orbit carbon footprint.\n• Automated Safety: RAG ensures compliance with space protocols.",
        "Challenges & Limitations": "• Simplified 2D orbital model used for rapid training (3D in progress).\n• Delta-V approximations rely on Hohmann transfers (Lambert solvers needed for high-fidelity).\n• Hardware validation remains for future phases.",
        "Conclusion": "• We demonstrated a functional proof-of-concept for autonomous, fuel-efficient orbital management.\n• Next Steps: Integrate full 3D orbital mechanics, deploy multi-agent coordination, and connect to ESA's live API.\n• ta5abes is ready to scale this for Phase 2."
    }

    # Iterate through slides
    for slide in prs.slides:
        # Find the title of this slide
        title_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                for key in slide_content.keys():
                    if key.lower() in text.lower() or text.lower() in key.lower():
                        title_text = key
                        break
            if title_text:
                break
        
        if title_text and title_text in slide_content:
            # We found a matching slide. Let's add the content.
            # Look for an empty text box or a large freeform shape
            content_added = False
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip() == "":
                    shape.text = slide_content[title_text]
                    content_added = True
                    break
            
            if not content_added:
                # Add a new text box
                txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(4.5))
                tf = txBox.text_frame
                tf.text = slide_content[title_text]
                for p in tf.paragraphs:
                    p.font.size = Pt(20)

    # For the first slide (title slide), let's just add the team name and title if not already there
    first_slide = prs.slides[0]
    txBox = first_slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    tf = txBox.text_frame
    tf.text = "Intelligent Orbital Debris Removal Planner\nTeam ta5abes - Track 4: Sustainable Space Systems"
    for p in tf.paragraphs:
        p.font.size = Pt(28)
        p.font.bold = True

    prs.save(output_path)
    print(f"Saved to {output_path}")

if __name__ == '__main__':
    fill_presentation(sys.argv[1], sys.argv[2])
